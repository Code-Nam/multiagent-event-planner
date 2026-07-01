"""
generate-ppt.py — AGEVP Phase 2

Reads a JSON spec file (ppt_presentation shape) and produces a .pptx
presentation. The first slide uses the title-slide layout; subsequent
slides use a title-and-content layout with bullet points.

If templates/presentation.pptx exists it is loaded as a theme base
(slide master, layouts, branding). Pass --template to override explicitly.

Usage:
    python scripts/generate-ppt.py [--input doc-content/slides.json] [--template <path>] [--output output/]
"""

import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing dependency. Run: pip install -r scripts/requirements.txt", file=sys.stderr)
    sys.exit(1)

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"
DEFAULT_TEMPLATE = TEMPLATES_DIR / "presentation.pptx"


def slugify(title: str) -> str:
    """Convert a title string to a lowercase underscore slug."""
    return title.lower().replace(" ", "_")


def _get_layout(prs: Presentation, preferred_index: int, fallback_index: int) -> object:
    """
    Return a slide layout by index, falling back gracefully when the
    theme has fewer layouts than expected.
    """
    layouts = prs.slide_layouts
    if preferred_index < len(layouts):
        return layouts[preferred_index]
    if fallback_index < len(layouts):
        return layouts[fallback_index]
    return layouts[0]


def _base_presentation(template: Path | None) -> Presentation:
    if template and template.exists():
        prs = Presentation(str(template))
        xml_slides = prs.slides._sldIdLst
        for slide_id in list(xml_slides):
            xml_slides.remove(slide_id)
        return prs
    return Presentation()


def generate_ppt(spec_path: Path, output_dir: Path, template: Path | None = None) -> Path:
    """
    Build a .pptx presentation from a ppt_presentation JSON spec.

    Args:
        spec_path:  Path to the JSON spec file.
        output_dir: Directory where the .pptx file will be written.

    Returns:
        Absolute path of the created .pptx file.

    Raises:
        SystemExit: on any I/O or validation error.
    """
    try:
        raw = spec_path.read_text(encoding="utf-8")
    except OSError as exc:
        print(f"ERROR: cannot read {spec_path}: {exc}", file=sys.stderr)
        sys.exit(1)

    try:
        spec = json.loads(raw)
    except json.JSONDecodeError as exc:
        print(f"ERROR: invalid JSON in {spec_path}: {exc}", file=sys.stderr)
        sys.exit(1)

    if spec.get("type") != "ppt_presentation":
        print(
            f"ERROR: expected type 'ppt_presentation', got '{spec.get('type')}'",
            file=sys.stderr,
        )
        sys.exit(1)

    title: str = spec.get("title", "Presentation")
    slides: list = spec.get("slides", [])

    resolved_template = template if template is not None else DEFAULT_TEMPLATE
    prs = _base_presentation(resolved_template)

    # Slide layout indices for the default Office Theme:
    #   0 = Title Slide, 1 = Title and Content
    title_slide_layout = _get_layout(prs, preferred_index=0, fallback_index=0)
    content_layout = _get_layout(prs, preferred_index=1, fallback_index=0)

    # First slide: presentation title
    title_slide = prs.slides.add_slide(title_slide_layout)
    placeholders = title_slide.placeholders
    if len(placeholders) > 0:
        placeholders[0].text = title

    for slide_spec in slides:
        slide_title: str = slide_spec.get("title", "")
        bullets: list = slide_spec.get("bullets", [])

        slide = prs.slides.add_slide(content_layout)
        slide_placeholders = slide.placeholders

        if len(slide_placeholders) > 0 and slide_title:
            slide_placeholders[0].text = slide_title

        if len(slide_placeholders) > 1 and bullets:
            tf = slide_placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    para = tf.add_paragraph()
                    para.text = bullet

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"{slugify(title)}.pptx"

    try:
        prs.save(out_path)
    except OSError as exc:
        print(f"ERROR: cannot write {out_path}: {exc}", file=sys.stderr)
        sys.exit(1)

    return out_path.resolve()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate a .pptx presentation from a ppt_presentation JSON spec."
    )
    parser.add_argument(
        "--input",
        default="doc-content/slides.json",
        help="Path to the JSON spec file (default: doc-content/slides.json)",
    )
    parser.add_argument(
        "--template",
        default=None,
        help=f"Path to a .pptx template for branding (default: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument(
        "--output",
        default="output",
        help="Output directory (default: output/)",
    )
    args = parser.parse_args()

    result = generate_ppt(
        spec_path=Path(args.input),
        output_dir=Path(args.output),
        template=Path(args.template) if args.template else None,
    )
    print(result)
    sys.exit(0)
