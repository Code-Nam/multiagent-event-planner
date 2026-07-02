"""
tokenize_templates.py — AGEVP Phase 2 (one-time / idempotent)

Injects {{TOKEN}} markers into the branded AGEVP templates in place, replacing
the demo `[ ... ]` bracket text at the slots the generators fill. All fonts,
colours, and layout are preserved — only placeholder run text changes.

A one-time `<template>.bak` backup is written on first run. Re-running is a
no-op for any template that already contains tokens.

Token catalog
  Modele_AGEVP.docx (report)
    {{DOC_TITLE}}     cover title
    {{DOC_SUBTITLE}}  cover subtitle
    {{SEC_NUM}}       section number run   (clone template)
    {{SEC_TITLE}}     section heading run  (clone template)
    {{SEC_BODY}}      section body para    (clone template)
  Modele_AGEVP_Presentation.pptx (presentation)
    {{PRES_TITLE}} {{PRES_SUBTITLE}} {{PRES_DATE}}   cover slide
    {{SLIDE_TITLE}}                                  content slide title
    {{POINT1}} {{POINT1_DESC}} .. {{POINT3_DESC}}    content slide 3 key points

Usage:
    python scripts/tokenize_templates.py
"""

import sys
import shutil
from pathlib import Path

try:
    from docx import Document
    from pptx import Presentation
except ImportError:
    print("Missing dependency. Run: pip install -r scripts/requirements.txt", file=sys.stderr)
    sys.exit(1)

import token_fill as tf  # noqa: E402  (same dir on sys.path when run as a script)

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"
DOCX = TEMPLATES_DIR / "Modele_AGEVP.docx"
PPTX = TEMPLATES_DIR / "Modele_AGEVP_Presentation.pptx"


def _backup(path: Path) -> None:
    bak = path.with_suffix(path.suffix + ".bak")
    if not bak.exists():
        shutil.copy2(path, bak)
        print(f"backup  {bak.relative_to(Path.cwd())}")


def _already_tokenized_docx(doc) -> bool:
    return any("{{" in tf.para_text(p) for p in doc.element.body.iter(f"{{{tf.W}}}p"))


def _already_tokenized_pptx(prs) -> bool:
    return any("{{" in tf.shape_text(sh)
               for slide in prs.slides for sh in slide.shapes)


def tokenize_docx() -> str:
    if not DOCX.exists():
        return f"skip  {DOCX.name} (not found)"
    doc = Document(str(DOCX))
    if _already_tokenized_docx(doc):
        return f"skip  {DOCX.name} (already tokenized)"

    for p in doc.element.body.iter(f"{{{tf.W}}}p"):
        txt = tf.para_text(p)
        if txt == "[ Titre du document ]":
            tf.set_first_run_text(p, "{{DOC_TITLE}}")
        elif txt.startswith("[ Sous-titre"):
            tf.set_first_run_text(p, "{{DOC_SUBTITLE}}")
        elif txt.startswith("1") and "[ Titre de section ]" in txt:
            runs = p.findall(f"{{{tf.W}}}r")
            if len(runs) >= 2:
                runs[0].find(f"{{{tf.W}}}t").text = "{{SEC_NUM}}"
                runs[1].find(f"{{{tf.W}}}t").text = "{{SEC_TITLE}}"
        elif txt.startswith("[ Paragraphe d'introduction"):
            tf.set_first_run_text(p, "{{SEC_BODY}}")

    _backup(DOCX)
    doc.save(str(DOCX))
    return f"ok    {DOCX.name} tokenized"


# pptx: (shape name -> whole-run token) for a full-placeholder shape
_SLIDE0 = {"Text 3": "{{PRES_TITLE}}", "Text 4": "{{PRES_SUBTITLE}}"}
_SLIDE2 = {
    "Text 1": "{{SLIDE_TITLE}}",
    "Text 4": "{{POINT1}}", "Text 5": "{{POINT1_DESC}}",
    "Text 8": "{{POINT2}}", "Text 9": "{{POINT2_DESC}}",
    "Text 12": "{{POINT3}}", "Text 13": "{{POINT3_DESC}}",
}


def _find_slide_with_shapes(prs, names: set) -> "object | None":
    """Return the first slide containing every shape name in `names`, or None.

    Slides are located by their shape inventory rather than hard index so a
    reordered or extended template still tokenizes the right slides.
    """
    for slide in prs.slides:
        if names <= {sh.name for sh in slide.shapes}:
            return slide
    return None


def _tokenize_slide(slide, mapping: dict) -> None:
    for shape in slide.shapes:
        if shape.name in mapping and shape.has_text_frame:
            tf.set_shape_first_run(shape, mapping[shape.name])


def _tokenize_date_shape(slide) -> None:
    """Replace only the '[ Date ]' run on the cover, keeping the ' • Annam…' suffix."""
    for shape in slide.shapes:
        if shape.name != "Text 6" or not shape.has_text_frame:
            continue
        for t in shape.text_frame._txBody.iter(f"{{{tf.A}}}t"):
            if t.text and "[ Date ]" in t.text:
                t.text = t.text.replace("[ Date ]", "{{PRES_DATE}}")
                return


def tokenize_pptx() -> str:
    if not PPTX.exists():
        return f"skip  {PPTX.name} (not found)"
    prs = Presentation(str(PPTX))
    if _already_tokenized_pptx(prs):
        return f"skip  {PPTX.name} (already tokenized)"

    cover = _find_slide_with_shapes(prs, set(_SLIDE0))
    content = _find_slide_with_shapes(prs, set(_SLIDE2))
    if cover is None or content is None:
        return f"skip  {PPTX.name} (cover/content shapes not found — template layout changed?)"

    _tokenize_slide(cover, _SLIDE0)
    _tokenize_date_shape(cover)
    _tokenize_slide(content, _SLIDE2)

    _backup(PPTX)
    prs.save(str(PPTX))
    return f"ok    {PPTX.name} tokenized"


def main() -> None:
    print(tokenize_docx())
    print(tokenize_pptx())
    print("\nTemplates ready. Generators locate slots by {{TOKEN}}, not position.")


if __name__ == "__main__":
    main()
