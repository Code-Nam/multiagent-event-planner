"""
generate_ppt.py — AGEVP Phase 2

Reads a JSON spec file (ppt_presentation shape) and produces a .pptx
presentation by cloning the branded AGEVP template slides and filling
{{TOKEN}} placeholders with spec content. The cover and content template
slides are found by token text, not slide index or shape name.

Tokens (see scripts/tokenize_templates.py):
  {{PRES_TITLE}} {{PRES_SUBTITLE}} {{PRES_DATE}}          cover slide
  {{SLIDE_TITLE}} {{POINT1..3}} {{POINT1..3_DESC}}        content slide

Spec shape:
  {"type": "ppt_presentation", "title": str, "date"?: str, "event"?: str,
   "slides": [{"title": str, "bullets": [str, ...]}, ...]}

Usage:
    python scripts/generate_ppt.py <json-spec-path> [--template <path>] [--output <dir>]
"""

import sys
import copy
import argparse
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Missing dependency. Run: pip install -r scripts/requirements.txt", file=sys.stderr)
    sys.exit(1)

from token_fill import load_spec, slide_has_token, fill_slide_tokens

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"
DEFAULT_TEMPLATE = TEMPLATES_DIR / "Modele_AGEVP_Presentation.pptx"

_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _find_slide(prs, token: str) -> int:
    for i, slide in enumerate(prs.slides):
        if slide_has_token(slide, token):
            return i
    return -1


def _clone_slide(prs, src_idx: int):
    """Append a deep-copy of prs.slides[src_idx] to prs and return the new slide.

    Uses the private `_spTree` element — python-pptx has no public slide-copy
    API. Tested against python-pptx 0.6/1.0; re-check on upgrade.
    """
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    dst_tree = new_slide.shapes._spTree
    for elem in list(dst_tree):
        dst_tree.remove(elem)
    for elem in src.shapes._spTree:
        dst_tree.append(copy.deepcopy(elem))
    return new_slide


def _remove_slide(prs, idx: int) -> None:
    """Remove prs.slides[idx] and its relationship entry.

    Uses the private `_sldIdLst`/`_rels` internals — python-pptx has no public
    slide-delete API. Tested against python-pptx 0.6/1.0; re-check on upgrade.
    """
    slide_ids = list(prs.slides._sldIdLst)
    if idx >= len(slide_ids):
        return
    r_id = slide_ids[idx].get(f"{{{_REL}}}id")
    if r_id:
        try:
            prs.part._rels.pop(r_id)
        except KeyError:
            print(f"WARNING: slide relationship {r_id} not found — "
                  "python-pptx internals may have changed", file=sys.stderr)
    prs.slides._sldIdLst.remove(slide_ids[idx])


def generate_ppt(spec_path: Path, output_dir: Path, template: Path | None = None) -> Path:
    spec = load_spec(spec_path, "ppt_presentation")

    resolved_template = template if template is not None else DEFAULT_TEMPLATE
    if not (resolved_template and resolved_template.exists()):
        print(f"ERROR: template not found: {resolved_template}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(resolved_template))
    n_original = len(prs.slides)

    cover_idx = _find_slide(prs, "PRES_TITLE")
    content_idx = _find_slide(prs, "POINT1")
    if cover_idx < 0 or content_idx < 0:
        print(
            f"ERROR: {resolved_template.name} is not tokenized "
            "(missing {{PRES_TITLE}}/{{POINT1}}). Run scripts/tokenize_templates.py.",
            file=sys.stderr,
        )
        sys.exit(1)

    title = spec.get("title", "Présentation")
    date = spec.get("date", "")
    event = spec.get("event", "")
    slides_spec = spec.get("slides", [])

    cover = _clone_slide(prs, cover_idx)
    fill_slide_tokens(cover, {
        "PRES_TITLE": title,
        "PRES_SUBTITLE": event or title,
        "PRES_DATE": date,
    })

    # Content slides — 3 bullets per output slide using the 3-key-points layout.
    for slide_spec in slides_spec:
        slide_title = slide_spec.get("title", "")
        bullets = slide_spec.get("bullets", [])
        chunks = [bullets[i:i + 3] for i in range(0, max(len(bullets), 1), 3)]
        for ci, chunk in enumerate(chunks):
            s = _clone_slide(prs, content_idx)
            label = slide_title if ci == 0 else f"{slide_title} (suite {ci + 1})"
            tokens = {"SLIDE_TITLE": label}
            for pi in range(3):
                tokens[f"POINT{pi + 1}"] = chunk[pi] if pi < len(chunk) else ""
                tokens[f"POINT{pi + 1}_DESC"] = ""
            fill_slide_tokens(s, tokens)

    # Remove the original template slides (always at indices 0..n_original-1).
    for _ in range(n_original):
        _remove_slide(prs, 0)

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"{spec_path.stem}.pptx"
    try:
        prs.save(out_path)
    except OSError as exc:
        print(f"ERROR: cannot write {out_path}: {exc}", file=sys.stderr)
        sys.exit(1)
    return out_path.resolve()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate a .pptx presentation from a ppt_presentation JSON spec."
    )
    parser.add_argument("spec", help="Path to the JSON spec file")
    parser.add_argument(
        "--template",
        default=None,
        help=f"Path to a .pptx template for branding (default: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument("--output", default="output", help="Output directory (default: output/)")
    args = parser.parse_args()

    result = generate_ppt(
        spec_path=Path(args.spec),
        output_dir=Path(args.output),
        template=Path(args.template) if args.template else None,
    )
    print(result)
    sys.exit(0)
